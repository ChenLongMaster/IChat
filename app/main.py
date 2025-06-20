import os
import traceback
import glob
import httpx
import math
import json
from urllib.parse import urlparse

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse

from pydantic import BaseModel
from dotenv import load_dotenv
from huggingface_hub import InferenceClient
from requests.exceptions import HTTPError

from qdrant_client import QdrantClient
from qdrant_client.models import Filter, FieldCondition, MatchValue

from pptx import Presentation

from llama_index.vector_stores.qdrant import QdrantVectorStore
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import (
    SimpleDirectoryReader,
    StorageContext,
    VectorStoreIndex,
    Settings
)
from llama_index.readers.web import BeautifulSoupWebReader
from llama_index.core.schema import Document

# Ragas
from datasets import Dataset
from ragas import evaluate
from ragas.metrics import faithfulness, answer_relevancy, context_precision, context_recall

# Ragas model wrappers
from ragas.llms import LangchainLLMWrapper
from ragas.embeddings import LangchainEmbeddingsWrapper

# LLM + Embeddings
from langchain_openai import ChatOpenAI
from ragas.llms import LangchainLLMWrapper
from langchain_community.embeddings import HuggingFaceEmbeddings

# === Load ENV ===
load_dotenv()


llm = LangchainLLMWrapper(ChatOpenAI(model="gpt-4o"))
wrapped_emb = LangchainEmbeddingsWrapper(
    HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
)

# === Constants ===
USE_GPU = os.getenv("USE_GPU", "False").lower() == "true"
DEFAULT_PROMPT = "You are a helpful assistant."
UPLOAD_FOLDER = r"C:\\IChat.Sources\\Upload"
PROMPT_FOLDER = r"C:\\IChat.Sources\\Instruction"
GENERIC_SETTING_FOLDER = r"C:\\IChat.Sources\\GeneralSetting"


qdrant_host = os.getenv("QDRANT_HOST", "localhost")  # default for local dev

# Load Power Point files
def load_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text.strip())
    return "\n".join([t for t in text_runs if t])

# === Models ===
class PromptRequest(BaseModel):
    tenant_id: str
    user_prompt: str

class CrawlRequest(BaseModel):
    tenant_id: str
    vector_id: str
    url: str

class TrainRequest(BaseModel):
    tenant_id: str
class DeleteVectorRequest(BaseModel):
    tenant_id: str
    fileNames: list[str]

# === Helpers ===
def load_system_prompt(data_dir: str, tenant_id: str) -> str:
    prompt_path = os.path.join(data_dir, f"{tenant_id}_prompt.txt")
    if os.path.exists(prompt_path):
        with open(prompt_path, "r", encoding="utf-8") as f:
            return f.read()
    return DEFAULT_PROMPT

def get_storage_paths(tenant_id: str):
    return (
        f"storage/{tenant_id}",
        f"ROOT_FOLDER\\{tenant_id}"
    )

# === Main app function ===
def main_app():
    app = FastAPI()

    app.add_middleware(
        CORSMiddleware,
        allow_origins=["*"],
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )

    Settings.embed_model = HuggingFaceEmbedding(model_name="sentence-transformers/all-MiniLM-L6-v2")

    def evaluate_rag_result(question: str, answer: str, contexts: list[str], reference: str | None = None) -> dict:
        try:
            entry = {
                "question": question,
                "answer": answer,
                "contexts": contexts
            }
            if reference:
                entry["reference"] = reference

            dataset = Dataset.from_list([entry])

            metrics = [faithfulness, answer_relevancy]
            if reference:
                metrics += [context_precision, context_recall]

            result = evaluate(
            result = evaluate(
                dataset,
                metrics=metrics,
                llm=llm,
                embeddings=wrapped_emb,
            )

            # ✅ Brute-force string parsing
            raw_str = str(result)  # Looks like: {'faithfulness': 0.9, 'answer_relevancy': 0.8}
            print("🧨 Raw string result:", raw_str)

            # Turn into real JSON-like dict
            import ast
            safe_result = ast.literal_eval(raw_str)

            return safe_result

        except Exception as e:
            print("⚠️ RAG evaluation failed:", e)
            return {"error": str(e)}





    @app.post("/bot")
    async def get_bot_response(req: PromptRequest):
        prompt_folder = os.path.join(PROMPT_FOLDER, req.tenant_id)
        system_prompt = DEFAULT_PROMPT
        temperature = 0.5
        max_tokens = 500

        if os.path.exists(prompt_folder):
            prompt_files = sorted(glob.glob(os.path.join(prompt_folder, "**", "*.txt"), recursive=True))

            if prompt_files:
                prompt_chunks = []
                for path in prompt_files:
                    with open(path, "r", encoding="utf-8") as f:
                        prompt_chunks.append(f.read())
                    print(f"Loaded prompt: {os.path.basename(path)}")

                system_prompt = "\n\n".join(prompt_chunks)
            else:
                print(f"No .txt files found in {prompt_folder}. Using default prompt.")
        else:
            print(f"⚠️ Folder not found: {prompt_folder}. Using default prompt.")

        setting_path = os.path.join(GENERIC_SETTING_FOLDER, req.tenant_id)
        if os.path.exists(setting_path):
            for dp, _, files in os.walk(setting_path):
                for file in files:
                    if not file.endswith(".txt"): continue
                    content = open(os.path.join(dp, file), "r", encoding="utf-8-sig").read().strip()
                    if content:
                        try:
                            config = json.loads(content)
                            temperature = config.get("Temperature", temperature)
                            max_tokens = config.get("MaxTokens", max_tokens)
                        except json.JSONDecodeError as e:
                            print(f"❌ JSON decode error in {file}: {e}")

        context = ""
        nodes = []
        try:
            qdrant_client = QdrantClient(host=qdrant_host, port=6333)
            collections = qdrant_client.get_collections().collections
            if req.tenant_id in [col.name for col in collections]:
                vector_store = QdrantVectorStore(client=qdrant_client, collection_name=req.tenant_id)
                index = VectorStoreIndex.from_vector_store(vector_store, storage_context=StorageContext.from_defaults(vector_store=vector_store))
                retriever = index.as_retriever(similarity_top_k=3)
                nodes = retriever.retrieve(req.user_prompt)
                context = "\n\n".join([n.get_content() for n in nodes])
        except Exception as e:
            print(f"❌ Qdrant access failed: {str(e)}")

        user_prompt = f"{req.user_prompt}\n\nContext:\n{context}"

        try:
            # Choose model dynamically (env var, tenant config, etc.)
            use_huggingface = False

            if use_huggingface:
                print("📡 Calling Hugging Face model...")
                result = await call_huggingface_model(user_prompt, system_prompt, temperature, max_tokens)
            else:
                print("⚙️ Calling Ollama self-hosted model...")
                result = await call_local_ollama_model(user_prompt, system_prompt, temperature, max_tokens)

            rag_scores = evaluate_rag_result(
                question=req.user_prompt,
                answer=result,
                contexts=[n.get_content() for n in nodes]
            )

            response = {
                "response": result,
                "rag_evaluation": rag_scores
            }

            # 🖨️ Print to console for live demo
            print("\n📤 Final Output to User:")
            print(json.dumps(response, indent=2, ensure_ascii=False))
            return response
        
        except Exception as e:
            print("❌ Model call failed:", traceback.format_exc())
            raise HTTPException(status_code=500, detail="Model call failed.")
        
    async def call_local_ollama_model(prompt: str, system_prompt: str, temperature: float, max_tokens: int):
        full_prompt = f"{system_prompt}\n\n{prompt}"
        print(full_prompt)
        payload = {
            "model": "llama3",
            "prompt": full_prompt,
            "stream": False
        }
        async with httpx.AsyncClient(timeout=360) as client:
            r = await client.post("https://ollama.s3corp.com.vn/api/generate", json=payload)
            r.raise_for_status()
            data = r.json()
            return data.get("response") or str(data)
        
    async def call_huggingface_model(prompt: str, system_prompt: str, temperature: float, max_tokens: int):
        try:
            client = InferenceClient(provider="cerebras", api_key = "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")
            response = client.chat.completions.create(
                model="meta-llama/Llama-3.3-70B-Instruct",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                temperature=temperature,
                max_tokens=max_tokens,
                top_p=0.7,
                stream=False
            )
            return response.choices[0].message["content"]

        except HTTPError as e:
            status = getattr(e.response, "status_code", "unknown")
            text = getattr(e.response, "text", str(e))
            if status == 429:
                raise HTTPException(status_code=429, detail="Hugging Face rate limit reached. Please try again later.")
            elif status == 401:
                raise HTTPException(status_code=401, detail="Invalid or missing Hugging Face API token.")
            elif status == 403:
                raise HTTPException(status_code=403, detail="Access to this model is restricted or you’ve hit your token quota.")
            else:
                raise HTTPException(status_code=500, detail=f"Inference API error: {text}")

    @app.post("/train-rag")
    async def train_rag(train_request: TrainRequest):
        try:
            tenant_id = train_request.tenant_id
            tenant_data_dir = os.path.join(UPLOAD_FOLDER, tenant_id)

            if not os.path.exists(tenant_data_dir):
                raise HTTPException(status_code=400, detail=f"Folder '{tenant_data_dir}' does not exist.")

            input_files = []

            for root, _, files in os.walk(tenant_data_dir):
                for f in files:
                    full_path = os.path.join(root, f)
                    input_files.append(full_path)

            if not input_files:
                raise HTTPException(status_code=400, detail=f"No valid files found to process in '{tenant_data_dir}'.")

            for file_path in input_files:
                print(f" - {file_path}")

            loaded_docs = []

            for file_path in input_files:
                ext = os.path.splitext(file_path)[1].lower()

                if ext == ".pptx":
                    try:
                        text = load_pptx(file_path)
                        doc = Document(text=text, metadata={"file_path": file_path})
                        loaded_docs.append(doc)
                        print(f"✅ Loaded PPTX: {file_path}")
                    except Exception as e:
                        print(f"⚠️ Failed to load PPTX: {file_path} -> {e}")
                else:
                    docs = SimpleDirectoryReader(input_files=[file_path]).load_data()
                    loaded_docs.extend(docs)

            documents = []

            for doc in loaded_docs:
                file_name = os.path.basename(doc.metadata.get("file_path", ""))
                doc.metadata["file_name"] = file_name
                print(f"📄 Indexed document with file_name: {file_name}")
                documents.append(doc)

            client = QdrantClient(host="localhost", port=6333)
            vector_store = QdrantVectorStore(client=client, collection_name=tenant_id)
            storage_context = StorageContext.from_defaults(vector_store=vector_store)
            VectorStoreIndex.from_documents(documents, storage_context=storage_context)

            return {"message": f"RAG model for tenant '{tenant_id}' trained and updated successfully."}

        except Exception as e:
            print("❌ ERROR OCCURRED:")
            print(traceback.format_exc())
            raise HTTPException(status_code=500, detail={
                "error": str(e),
                "trace": traceback.format_exc()
            })

    @app.post("/delete-vectors")
    async def delete_vectors(request: DeleteVectorRequest):
        try:
            qdrant_client = QdrantClient(host="localhost", port=6333)

            # Confirm collection exists
            collections = qdrant_client.get_collections().collections
            if request.tenant_id not in [col.name for col in collections]:
                raise HTTPException(status_code=404, detail=f"Collection '{request.tenant_id}' not found in Qdrant.")

            deleted = 0
            for file_name in request.fileNames:
                try:
                    qdrant_client.delete(
                        collection_name=request.tenant_id,
                        points_selector=Filter(
                            must=[
                                FieldCondition(
                                    key="file_name",
                                    match=MatchValue(value=file_name)
                                )
                            ]
                        )
                    )
                    print(f"🗑️ Deleted all vectors with file_name = {file_name}")
                    deleted += 1
                except Exception as ve:
                    print(f"⚠️ Failed to delete file {file_name}: {str(ve)}")

            return {"message": f"Deleted {deleted} out of {len(request.fileNames)} file(s)."}

        except Exception as e:
            print(traceback.format_exc())
            raise HTTPException(status_code=500, detail={
                "error": str(e),
                "trace": traceback.format_exc()
            })


    @app.post("/crawl-job")
    async def crawl_job(crawl_request: CrawlRequest):
        url = crawl_request.url
        tenant_id = crawl_request.tenant_id
        vector_id = crawl_request.vector_id

        try:
            domain = urlparse(url).netloc.replace("www.", "")
            loader = BeautifulSoupWebReader()
            documents = loader.load_data(urls=[url])

            if not documents:
                raise HTTPException(status_code=404, detail="No content found on the page.")

            crawl_dir = os.path.join(UPLOAD_FOLDER, tenant_id, "crawl")
            os.makedirs(crawl_dir, exist_ok=True)

            for doc in documents:
                doc.metadata = {"tenant_id": tenant_id, "source_type": "web_crawl", "source_origin": url, "vector_id": vector_id}
                file_path = os.path.join(crawl_dir, f"{domain}_{vector_id}.txt")
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(f"URL: {url}\n\n" + "\n".join(line.strip() for line in doc.text.splitlines() if line.strip()))

            return {"message": "Content fetched and saved.", "url": url, "vector_id": vector_id, "document_count": len(documents)}
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to fetch and save URL: {str(e)}")

    return app

# === Run App ===
app = main_app()
